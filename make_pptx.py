from pptx import Presentation
from pptx.util import Inches
import os

extension = '.PNG'
# Get a list of PNG files in the current directory
png_files = [file for file in os.listdir() if file.endswith(extension)]

# Sort the list of filenames numerically
png_files = sorted(png_files, key=lambda x: int(''.join(filter(str.isdigit, x))))

print("PNG files found in the directory:", png_files)

if not png_files:
    print("No PNG files found in the directory.")
    exit()

# Get the dimensions of the first  image
first_image = png_files[0]
image_width_px, image_height_px = 0, 0

try:
    from PIL import Image
    with Image.open(first_image) as img:
        image_width_px, image_height_px = img.size
except Exception as e:
    print("Error:", e)
    exit()

# Convert image dimensions from pixels to inches
image_width_inches = image_width_px / 96  # Assuming 96 pixels per inch (standard screen resolution)
image_height_inches = image_height_px / 96

# Create a presentation object
presentation = Presentation()

# Set slide size based on the dimensions of the first image
presentation.slide_width = Inches(image_width_inches)
presentation.slide_height = Inches(image_height_inches)

# Add each PNG image to individual slides in the presentation
for png_file in png_files:
    # Add a blank slide to the presentation
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Using layout index 5
    
    # Add the PNG image to the slide
    left = top = Inches(0)
    slide.shapes.add_picture(png_file, left, top, width=Inches(image_width_inches), height=Inches(image_height_inches))

# Save the presentation to a file
presentation.save("Created_Presentaion.pptx")
print("Presentation created successfully!")
